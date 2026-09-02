"""
读 PPT 内容
"""
from pptx import Presentation
import sys

PPT = r"C:\Users\34464\.minimax\v2\assets\2026\09\02\18-00-04-327-asset_20260902-180004-327_22a3fc393641_2094d974-RiskLens-答辩版-可信法律风险分析-截图清晰版.pptx"

p = Presentation(PPT)
print(f"Total slides: {len(p.slides)}")
print(f"Slide size: {p.slide_width} x {p.slide_height}")

for i, slide in enumerate(p.slides, 1):
    print(f"\n========== Slide {i} ==========")
    for sh in slide.shapes:
        if sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                t = "".join(r.text for r in para.runs)
                if t.strip():
                    print(f"  TEXT: {t.strip()[:200]}")
        elif sh.shape_type == 13:  # picture
            print(f"  [PICTURE: {sh.name}]")
        elif sh.has_table:
            tbl = sh.table
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                print(f"  TABLE: {' | '.join(cells)[:200]}")
        else:
            print(f"  SHAPE: {sh.shape_type} {sh.name}")
